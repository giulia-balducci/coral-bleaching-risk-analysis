from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette (matches NLP deck) ───────────────────────────────────────────────
DARK_TEAL = RGBColor(0x1B, 0x3D, 0x3D)
TEAL      = RGBColor(0x00, 0x83, 0x8A)
CORAL     = RGBColor(0xE0, 0x70, 0x50)
OFF_WHITE = RGBColor(0xF4, 0xF4, 0xF0)
DARK_NAVY = RGBColor(0x1A, 0x25, 0x33)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MID_GREY  = RGBColor(0x88, 0x88, 0x88)
LT_TEAL   = RGBColor(0xE8, 0xF4, 0xF4)
LT_GREY   = RGBColor(0xCC, 0xCC, 0xCC)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0.5)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def tx(slide, text, l, t, w, h,
       size=14, bold=False, italic=False,
       color=DARK_NAVY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = 'Calibri'
    return tb

def title(slide, text, dark_bg=False):
    c = WHITE if dark_bg else DARK_NAVY
    tx(slide, text, Inches(0.5), Inches(0.4), W - Inches(1.0), Inches(0.8),
       size=34, bold=True, color=c)
    rect(slide, Inches(0.5), Inches(1.25), W - Inches(1.0), Pt(3), fill=TEAL)

def caption(slide, text, y=Inches(1.4)):
    tx(slide, text, Inches(0.5), y, W - Inches(1.0), Inches(0.4),
       size=13, italic=True, color=TEAL)

def vbar(slide, x, y, h):
    rect(slide, x, y, Pt(3), h, fill=TEAL)

def bottom_box(slide, text):
    rect(slide, 0, H - Inches(1.05), W, Inches(1.05), fill=DARK_TEAL)
    tx(slide, text, Inches(0.5), H - Inches(0.95), W - Inches(1.0), Inches(0.85),
       size=13, bold=True, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

def photo_placeholder(slide, l, t, w, h, label='[ Insert reef photo ]'):
    rect(slide, l, t, w, h,
         fill=RGBColor(0x2A, 0x50, 0x50),
         line=TEAL, lw=Pt(1))
    tx(slide, label, l, t + h/2 - Inches(0.3), w, Inches(0.6),
       size=12, italic=True, color=TEAL, align=PP_ALIGN.CENTER)

def coral_strip(slide):
    rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=CORAL)

def img_placeholder(slide, l, t, w, h, label='[ Insert chart / screenshot ]'):
    rect(slide, l, t, w, h, fill=LT_TEAL, line=LT_GREY)
    tx(slide, label, l, t + h/2 - Inches(0.3), w, Inches(0.6),
       size=12, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)


# ── SLIDE 1: Cover ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, DARK_TEAL)

tx(s, 'Coral Bleaching Risk Analysis',
   Inches(0.6), Inches(1.8), Inches(6.4), Inches(1.2),
   size=42, bold=True, color=WHITE)
tx(s, 'Predicting bleaching events in the Atlantic Ocean',
   Inches(0.6), Inches(3.0), Inches(6.4), Inches(0.6),
   size=19, italic=True, color=TEAL)
rect(s, Inches(0.6), Inches(3.75), Inches(3.2), Pt(3), fill=TEAL)
tx(s, 'Giulia Balducci',
   Inches(0.6), Inches(3.95), Inches(6.0), Inches(0.4),
   size=16, bold=True, color=WHITE)
tx(s, 'AllWomen Applied AI Bootcamp  ·  April 2026',
   Inches(0.6), Inches(4.4), Inches(6.0), Inches(0.4),
   size=13, color=MID_GREY)

photo_placeholder(s, Inches(7.5), Inches(0.75), Inches(5.3), Inches(6.0))
coral_strip(s)


# ── SLIDE 2: The Problem ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, OFF_WHITE)
title(s, 'The Problem')

vbar(s, Inches(0.5), Inches(1.55), Inches(4.3))
tx(s, '50% of coral reefs lost since the 1980s.\n\nRising sea surface temperatures are the primary driver of bleaching — a thermal stress response that can kill entire reef ecosystems.\n\nI have witnessed this firsthand: 1,600 dives, over 1,200 in the Coral Triangle.',
   Inches(0.75), Inches(1.6), Inches(5.4), Inches(3.8),
   size=14, color=DARK_NAVY)

vbar(s, W/2, Inches(1.55), Inches(4.3))
tx(s, 'The Opportunity',
   Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.45),
   size=17, bold=True, color=DARK_NAVY)
tx(s, '→  Can we predict bleaching events before they happen?\n\n→  Which environmental variables drive bleaching risk?\n\n→  Can a regional model outperform a global one?',
   Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.2),
   size=14, color=DARK_NAVY)

bottom_box(s, '"Missing a bleaching event carries higher ecological cost than a false alarm."')


# ── SLIDE 3: Dataset ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, OFF_WHITE)
title(s, '40 years. 93 countries. 41,361 observations.')
caption(s, 'Van Woesik & Kratochwill (2022)  ·  Global Coral Bleaching Database  ·  BCO-DMO')

blocks = [
    ('1.', 'Dataset',          '62 variables: SST, DHW, depth,\nturbidity, cyclone frequency,\nlatitude, longitude, substrate'),
    ('2.', 'Temporal coverage','1980–2020  ·  40 years\nMajor bleaching events:\n1998, 2005, 2010, 2016'),
    ('3.', 'Target variable',  'Bleaching_Binary: none / bleaching\nGlobal: 82% none  ·  18% bleaching'),
    ('4.', 'Class imbalance',  'Handled with SMOTE\nAtlantic subset: 68% none  ·  32% bleaching\n(nearly twice as balanced as global)'),
]
positions = [
    (Inches(0.5),  Inches(2.05)),
    (Inches(7.0),  Inches(2.05)),
    (Inches(0.5),  Inches(4.6)),
    (Inches(7.0),  Inches(4.6)),
]
for (num, hd, body), (lx, ty) in zip(blocks, positions):
    rect(s, lx, ty, Inches(0.55), Inches(0.55), fill=TEAL)
    tx(s, num, lx + Inches(0.05), ty - Inches(0.02), Inches(0.45), Inches(0.6),
       size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    vbar(s, lx + Inches(0.65), ty, Inches(1.5))
    tx(s, hd, lx + Inches(0.8), ty, Inches(5.4), Inches(0.45),
       size=15, bold=True, color=DARK_NAVY)
    tx(s, body, lx + Inches(0.8), ty + Inches(0.5), Inches(5.4), Inches(1.2),
       size=13, color=DARK_NAVY)


# ── SLIDE 4: Approach ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, OFF_WHITE)
title(s, 'From global to regional — and why it matters')

steps = [
    ('Clean &\nfeature select', TEAL),
    ('Global model\n(RF + SMOTE)',  TEAL),
    ('Error analysis\n& SHAP',      TEAL),
    ('Atlantic\nregional model',    CORAL),
]
for i, (label, col) in enumerate(steps):
    lx = Inches(0.5 + i * 3.1)
    rect(s, lx, Inches(1.65), Inches(2.75), Inches(0.95), fill=col)
    tx(s, label, lx, Inches(1.7), Inches(2.75), Inches(0.85),
       size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        tx(s, '→', lx + Inches(2.75), Inches(1.95), Inches(0.35), Inches(0.45),
           size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

vbar(s, W/2, Inches(2.85), Inches(3.9))

tx(s, 'Why Atlantic?',
   Inches(0.5), Inches(2.9), Inches(5.8), Inches(0.45),
   size=15, bold=True, color=DARK_NAVY)
tx(s, 'Global model showed structural geographic bias:\nAtlantic FN rate 0.19 vs Pacific 0.74.\n\nRoot cause: Pacific major events (2016, 2020, 2022, 2024)\nfall outside the training window. Atlantic events (1998, 2005)\nare within training — model learns the right patterns.',
   Inches(0.5), Inches(3.4), Inches(5.8), Inches(3.1),
   size=13, color=DARK_NAVY)

tx(s, 'Metrics',
   Inches(7.0), Inches(2.9), Inches(5.8), Inches(0.45),
   size=15, bold=True, color=DARK_NAVY)
tx(s, 'Primary metric: Recall\nMissing a bleaching event > false alarm\n\nSecondary: Macro F1\n\nTemporal split:\ntrain 1980–2009  ·  val 2010–2014  ·  test 2015–2019',
   Inches(7.0), Inches(3.4), Inches(5.8), Inches(3.1),
   size=13, color=DARK_NAVY)


# ── SLIDE 5: Results ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, OFF_WHITE)
title(s, 'Atlantic recall: 0.78. Global: 0.64.')

rows = [
    ('Model',                            'Recall',  'Macro F1', 'header'),
    ('Global RF + SMOTE (baseline)',     '0.64',    '0.73',     'normal'),
    ('Atlantic vanilla',                 '0.73',    '0.72',     'normal'),
    ('Atlantic tuned',                   '0.82',    '0.67',     'normal'),
    ('Atlantic — no substrate  ✓ final', '0.78',    '0.69',     'highlight'),
]
col_x = [Inches(1.0), Inches(8.3), Inches(10.9)]
col_w = [Inches(7.0), Inches(2.3), Inches(2.2)]
rh = Inches(0.62)

for i, (c0, c1, c2, style) in enumerate(rows):
    ty = Inches(1.55) + i * rh
    bg_col = DARK_NAVY if style == 'header' else (TEAL if style == 'highlight' else OFF_WHITE)
    fg_col = WHITE if style in ('header', 'highlight') else DARK_NAVY
    line_col = LT_GREY if style == 'normal' else None
    rect(s, Inches(1.0), ty, Inches(12.1), rh,
         fill=bg_col, line=line_col, lw=Pt(0.5))
    for cx, cw, txt in zip(col_x, col_w, (c0, c1, c2)):
        tx(s, txt, cx + Inches(0.1), ty + Inches(0.12), cw, rh - Inches(0.1),
           size=13, bold=(style != 'normal'), color=fg_col)

bottom_box(s, '"Correctly identifies 8 out of 10 bleaching events using temperature, thermal stress accumulation, and seasonality."')


# ── SLIDE 6: SHAP ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, OFF_WHITE)
title(s, 'The model learns thermal physics, not survey artefacts')
caption(s, 'SHAP bar plot — with Substrate_Name (left) vs without (right)')

vbar(s, W/2, Inches(1.55), Inches(4.7))

# Left panel
tx(s, 'With Substrate_Name',
   Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.45),
   size=14, bold=True, color=DARK_NAVY)
img_placeholder(s, Inches(0.5), Inches(2.1), Inches(5.8), Inches(3.3),
                '[ Insert SHAP bar chart WITH substrate ]')
tx(s, '① Not_Recorded dominates at 0.10  (3× TSA_DHW)\n   → model learned survey methodology, not ecology',
   Inches(0.5), Inches(5.5), Inches(5.8), Inches(0.65),
   size=12, bold=True, color=CORAL)

# Right panel
tx(s, 'Without Substrate_Name  ✓ final',
   Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.45),
   size=14, bold=True, color=DARK_NAVY)
img_placeholder(s, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.3),
                '[ Insert SHAP bar chart WITHOUT substrate ]')
tx(s, '① Temperature_Mean 0.050   ④ TSA_DHW\n   → thermal physics. Recall cost: −3.8pp only.',
   Inches(7.0), Inches(5.5), Inches(5.8), Inches(0.65),
   size=12, bold=True, color=TEAL)


# ── SLIDE 7: Error analysis + Limits ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, OFF_WHITE)
title(s, 'What the model gets wrong — and why')

vbar(s, W/2, Inches(1.55), Inches(4.7))

# Left: stats
tx(s, 'Error analysis — validation set',
   Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.45),
   size=15, bold=True, color=DARK_NAVY)

stats = [
    ('0.22', 'Miss rate (FN)',        'vs 0.39 global  ↓'),
    ('0.31', 'False alarm rate (FP)', 'vs 0.34 global  ↓'),
    ('43%',  'FNs fall in 2010',      'post-El Niño anomaly dominates'),
]
for i, (val, label, note) in enumerate(stats):
    ty = Inches(2.2) + i * Inches(1.1)
    tx(s, val,   Inches(0.5), ty, Inches(1.2), Inches(0.7),
       size=30, bold=True, color=TEAL)
    tx(s, label, Inches(1.8), ty, Inches(4.0), Inches(0.4),
       size=13, bold=True, color=DARK_NAVY)
    tx(s, note,  Inches(1.8), ty + Inches(0.42), Inches(4.0), Inches(0.35),
       size=11, italic=True, color=MID_GREY)

# Right: timeline
tx(s, 'Training window vs bleaching events',
   Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.45),
   size=15, bold=True, color=DARK_NAVY)

# Timeline bar
rect(s, Inches(7.0), Inches(2.6), Inches(5.8), Inches(0.07), fill=MID_GREY)
events = [
    (0.00, '1980', '← train\nstart',     MID_GREY),
    (0.38, '1998', '1998 ✓\nin train',   TEAL),
    (0.63, '2005', '2005 ✓\nin train',   TEAL),
    (0.75, '2010', '2010\n43% FNs',      CORAL),
    (1.00, '2019', '← test\nend',        MID_GREY),
]
for frac, year, label, col in events:
    lx = Inches(7.0) + frac * Inches(5.8)
    rect(s, lx - Inches(0.07), Inches(2.47), Inches(0.14), Inches(0.32), fill=col)
    tx(s, year,  lx - Inches(0.4), Inches(2.78), Inches(0.8), Inches(0.35),
       size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    tx(s, label, lx - Inches(0.65), Inches(3.12), Inches(1.3), Inches(0.65),
       size=9, color=col, align=PP_ALIGN.CENTER)

tx(s, 'Pacific gap (2016–2024)\noutside training window',
   Inches(7.0), Inches(4.1), Inches(5.8), Inches(0.65),
   size=12, italic=True, color=MID_GREY)

bottom_box(s, 'Test set (2015–2019) not yet touched. All metrics above are on the validation set.')


# ── SLIDE 8: Demo + Future work ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, OFF_WHITE)
title(s, 'Atlantic Bleaching Risk — live demo')

vbar(s, W/2, Inches(1.55), Inches(5.2))

# Left: app
tx(s, 'Interactive Bleaching Risk App',
   Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.45),
   size=15, bold=True, color=DARK_NAVY)
img_placeholder(s, Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.1),
                '[ Insert app screenshot or GIF ]')
tx(s, 'Click reef site  →  sliders auto-fill  →  predict  →  bleaching probability',
   Inches(0.5), Inches(6.3), Inches(5.8), Inches(0.45),
   size=11, italic=True, color=TEAL)

# Right: future work
tx(s, 'What comes next',
   Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.45),
   size=15, bold=True, color=DARK_NAVY)

future = [
    ('Pacific model',        'Post-2020 data needed to capture\n2016, 2020, 2022, 2024 events'),
    ('LSTM for SST',         'Time series imputation\n& temperature forecasting'),
    ('Early warning system', 'Couple bleaching classifier\nwith DHW forecast model'),
]
for i, (hd, body) in enumerate(future):
    ty = Inches(2.2) + i * Inches(1.55)
    vbar(s, Inches(7.0), ty, Inches(1.2))
    tx(s, hd,   Inches(7.25), ty, Inches(5.5), Inches(0.45),
       size=14, bold=True, color=DARK_NAVY)
    tx(s, body, Inches(7.25), ty + Inches(0.48), Inches(5.5), Inches(0.85),
       size=12, color=MID_GREY)


# ── SLIDE 9: Thank you ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, DARK_TEAL)

tx(s, 'Thank you',
   Inches(0.6), Inches(1.8), Inches(6.2), Inches(1.1),
   size=50, bold=True, color=WHITE)
tx(s, 'Questions?',
   Inches(0.6), Inches(2.95), Inches(6.2), Inches(0.65),
   size=22, italic=True, color=TEAL)
rect(s, Inches(0.6), Inches(3.75), Inches(3.5), Pt(3), fill=TEAL)
tx(s, 'Giulia Balducci',
   Inches(0.6), Inches(3.95), Inches(6.2), Inches(0.45),
   size=16, bold=True, color=WHITE)
tx(s, 'linkedin.com/in/giuliabalducci   ·   github.com/giulia-balducci',
   Inches(0.6), Inches(4.45), Inches(6.2), Inches(0.45),
   size=13, color=MID_GREY)

photo_placeholder(s, Inches(7.5), Inches(0.75), Inches(5.3), Inches(5.95))
coral_strip(s)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), 'coral_bleaching_slides.pptx')
prs.save(out_path)
print(f"Saved: {out_path}")
